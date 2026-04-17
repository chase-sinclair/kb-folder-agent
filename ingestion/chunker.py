import csv
import re
from dataclasses import dataclass, field
from pathlib import Path


class UnsupportedFileTypeError(Exception):
    pass


@dataclass
class ChunkResult:
    content: str
    chunk_index: int
    chunk_type: str = "text"
    metadata: dict = field(default_factory=dict)


# ---------------------------------------------------------------------------
# Token helpers
# ---------------------------------------------------------------------------

def estimate_tokens(text: str) -> int:
    return int(len(text.split()) / 0.75)


def split_into_chunks(text: str, target_tokens: int, overlap_tokens: int) -> list[str]:
    words = text.split()
    target_words = int(target_tokens * 0.75)
    overlap_words = int(overlap_tokens * 0.75)

    if len(words) <= target_words:
        return [text] if words else []

    chunks: list[str] = []
    start = 0
    while start < len(words):
        end = min(start + target_words, len(words))
        chunks.append(" ".join(words[start:end]))
        if end == len(words):
            break
        start = end - overlap_words
    return chunks


# ---------------------------------------------------------------------------
# Router
# ---------------------------------------------------------------------------

async def chunk_file(file_path: str) -> list[ChunkResult]:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return chunk_pdf(file_path)
    if ext == ".docx":
        return chunk_docx(file_path)
    if ext in {".md", ".txt"}:
        return chunk_markdown(file_path)
    if ext in {".xlsx", ".csv"}:
        return chunk_spreadsheet(file_path)
    if ext in {".py", ".js", ".ts", ".go", ".rs"}:
        return chunk_code(file_path)
    if ext == ".pptx":
        return chunk_pptx(file_path)
    if ext == ".eml":
        return chunk_email(file_path)
    if ext == ".html":
        return chunk_html(file_path)
    raise UnsupportedFileTypeError(f"Unsupported file type: {ext!r}")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _table_to_markdown(table: list[list]) -> str:
    if not table:
        return ""
    rows = [[str(cell or "") for cell in row] for row in table]
    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
    return "\n".join(filter(None, [header, separator, body]))


def chunk_pdf(file_path: str) -> list[ChunkResult]:
    import pdfplumber

    results: list[ChunkResult] = []
    idx = 0

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            # Tables first
            tables = page.extract_tables() or []
            for table in tables:
                md = _table_to_markdown(table)
                if md.strip():
                    results.append(ChunkResult(
                        content=md,
                        chunk_index=idx,
                        chunk_type="table",
                        metadata={"page_number": page_num},
                    ))
                    idx += 1

            # Text (excluding table bounding boxes)
            text = page.extract_text() or ""
            if not text.strip():
                continue

            for chunk_text in split_into_chunks(text, target_tokens=600, overlap_tokens=75):
                if chunk_text.strip():
                    results.append(ChunkResult(
                        content=chunk_text,
                        chunk_index=idx,
                        chunk_type="text",
                        metadata={"page_number": page_num},
                    ))
                    idx += 1

    return results


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def chunk_docx(file_path: str) -> list[ChunkResult]:
    from docx import Document

    doc = Document(file_path)
    results: list[ChunkResult] = []
    idx = 0

    current_heading = ""
    current_paragraphs: list[str] = []

    def flush_section(heading: str, paragraphs: list[str]) -> None:
        nonlocal idx
        text = "\n".join(paragraphs).strip()
        if not text:
            return
        for chunk_text in split_into_chunks(text, target_tokens=600, overlap_tokens=75):
            if chunk_text.strip():
                results.append(ChunkResult(
                    content=chunk_text,
                    chunk_index=idx,
                    chunk_type="text",
                    metadata={"section_heading": heading},
                ))
                idx += 1

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading"):
            flush_section(current_heading, current_paragraphs)
            current_heading = text
            current_paragraphs = []
        else:
            current_paragraphs.append(text)

    flush_section(current_heading, current_paragraphs)
    return results


# ---------------------------------------------------------------------------
# Markdown / plain text
# ---------------------------------------------------------------------------

def chunk_markdown(file_path: str) -> list[ChunkResult]:
    text = Path(file_path).read_text(encoding="utf-8")
    results: list[ChunkResult] = []
    idx = 0

    # Split into raw blocks on blank lines, preserving fenced code blocks
    raw_blocks: list[tuple[str, bool]] = []  # (content, is_code)
    in_fence = False
    fence_buf: list[str] = []
    text_buf: list[str] = []

    for line in text.splitlines():
        if line.strip().startswith("```"):
            if in_fence:
                fence_buf.append(line)
                raw_blocks.append(("\n".join(fence_buf), True))
                fence_buf = []
                in_fence = False
            else:
                if text_buf:
                    raw_blocks.append(("\n".join(text_buf), False))
                    text_buf = []
                fence_buf = [line]
                in_fence = True
        elif in_fence:
            fence_buf.append(line)
        else:
            if line.strip() == "" and text_buf:
                raw_blocks.append(("\n".join(text_buf), False))
                text_buf = []
            elif line.strip():
                text_buf.append(line)

    if fence_buf:
        raw_blocks.append(("\n".join(fence_buf), True))
    if text_buf:
        raw_blocks.append(("\n".join(text_buf), False))

    # Merge small non-code paragraphs with next
    merged: list[tuple[str, bool]] = []
    i = 0
    while i < len(raw_blocks):
        content, is_code = raw_blocks[i]
        if (
            not is_code
            and len(content.split()) < 50
            and i + 1 < len(raw_blocks)
            and not raw_blocks[i + 1][1]
        ):
            merged.append((content + "\n\n" + raw_blocks[i + 1][0], False))
            i += 2
        else:
            merged.append((content, is_code))
            i += 1

    para_idx = 0
    for content, is_code in merged:
        if not content.strip():
            continue
        results.append(ChunkResult(
            content=content,
            chunk_index=idx,
            chunk_type="code" if is_code else "text",
            metadata={"paragraph_index": para_idx},
        ))
        idx += 1
        para_idx += 1

    return results


# ---------------------------------------------------------------------------
# Spreadsheet
# ---------------------------------------------------------------------------

def _rows_to_markdown(headers: list[str], rows: list[list]) -> str:
    header = "| " + " | ".join(str(h) for h in headers) + " |"
    separator = "| " + " | ".join("---" for _ in headers) + " |"
    body = "\n".join(
        "| " + " | ".join(str(cell) for cell in row) + " |"
        for row in rows
    )
    return "\n".join(filter(None, [header, separator, body]))


def chunk_spreadsheet(file_path: str) -> list[ChunkResult]:
    ext = Path(file_path).suffix.lower()
    results: list[ChunkResult] = []
    idx = 0

    def emit_row_groups(sheet_name: str, headers: list, all_rows: list) -> None:
        nonlocal idx
        for start in range(0, max(len(all_rows), 1), 50):
            group = all_rows[start: start + 50]
            if not group:
                continue
            end = start + len(group)
            md = _rows_to_markdown(headers, group)
            results.append(ChunkResult(
                content=md,
                chunk_index=idx,
                chunk_type="table",
                metadata={"sheet_name": sheet_name, "row_range": f"{start + 1}-{end}"},
            ))
            idx += 1

    if ext == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_rows = list(ws.iter_rows(values_only=True))
            if not all_rows:
                continue
            headers = [str(h) if h is not None else "" for h in all_rows[0]]
            data_rows = [[str(c) if c is not None else "" for c in r] for r in all_rows[1:]]
            emit_row_groups(sheet_name, headers, data_rows)
        wb.close()
    else:  # .csv
        with open(file_path, newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            all_rows = list(reader)
        if all_rows:
            headers = all_rows[0]
            data_rows = all_rows[1:]
            emit_row_groups(Path(file_path).stem, headers, data_rows)

    return results


# ---------------------------------------------------------------------------
# Code
# ---------------------------------------------------------------------------

_SPLIT_PATTERN = re.compile(r"(?=^(?:def |class |func |fn )\s*\w)", re.MULTILINE)


def chunk_code(file_path: str) -> list[ChunkResult]:
    text = Path(file_path).read_text(encoding="utf-8")
    language = Path(file_path).suffix.lstrip(".")
    results: list[ChunkResult] = []

    raw_blocks = _SPLIT_PATTERN.split(text)
    # Rejoin any leading non-definition content with the first block
    blocks = [b for b in raw_blocks if b.strip()]

    for idx, block in enumerate(blocks):
        results.append(ChunkResult(
            content=block.rstrip(),
            chunk_index=idx,
            chunk_type="code",
            metadata={"language": language},
        ))

    return results


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def chunk_pptx(file_path: str) -> list[ChunkResult]:
    from pptx import Presentation

    prs = Presentation(file_path)
    results: list[ChunkResult] = []
    idx = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type == 13:  # picture placeholder — skip
                continue
            if shape.name.lower().startswith("title") or (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            ):
                title = text
            else:
                body_parts.append(text)

        content = "\n".join(filter(None, [title] + body_parts)).strip()
        if not content:
            continue

        results.append(ChunkResult(
            content=content,
            chunk_index=idx,
            chunk_type="slide",
            metadata={"slide_number": slide_num, "slide_title": title},
        ))
        idx += 1

    return results


# ---------------------------------------------------------------------------
# Email (.eml)
# ---------------------------------------------------------------------------

def chunk_email(file_path: str) -> list[ChunkResult]:
    import email as _email
    import html.parser

    class _HTMLStripper(html.parser.HTMLParser):
        def __init__(self):
            super().__init__()
            self._parts: list[str] = []
        def handle_data(self, data: str) -> None:
            self._parts.append(data)
        def get_text(self) -> str:
            return " ".join(self._parts)

    raw = Path(file_path).read_bytes()
    msg = _email.message_from_bytes(raw)

    from_addr = msg.get("From", "")
    to_addr = msg.get("To", "")
    subject = msg.get("Subject", "")
    date = msg.get("Date", "")

    body_parts: list[str] = []
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                body_parts.append(part.get_payload(decode=True).decode(errors="replace"))
            elif ct == "text/html" and not body_parts:
                stripper = _HTMLStripper()
                stripper.feed(part.get_payload(decode=True).decode(errors="replace"))
                body_parts.append(stripper.get_text())
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            text = payload.decode(errors="replace")
            if msg.get_content_type() == "text/html":
                stripper = _HTMLStripper()
                stripper.feed(text)
                text = stripper.get_text()
            body_parts.append(text)

    body = "\n".join(body_parts).strip()
    content = f"From: {from_addr}\nTo: {to_addr}\nSubject: {subject}\nDate: {date}\n\n{body}"

    if not content.strip():
        return []

    return [ChunkResult(
        content=content,
        chunk_index=0,
        chunk_type="email",
        metadata={"from": from_addr, "to": to_addr, "subject": subject, "date": date},
    )]


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def chunk_html(file_path: str) -> list[ChunkResult]:
    from bs4 import BeautifulSoup

    raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")

    for tag in soup(["script", "style"]):
        tag.decompose()

    title = soup.title.string.strip() if soup.title and soup.title.string else ""
    text = soup.get_text(separator="\n")

    # Collapse blank lines into paragraph breaks then reuse markdown paragraph chunking
    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
    merged_text = "\n\n".join(paragraphs)

    results: list[ChunkResult] = []
    idx = 0
    for para_idx, para in enumerate(paragraphs):
        if not para:
            continue
        for chunk_text in split_into_chunks(para, target_tokens=600, overlap_tokens=75):
            if chunk_text.strip():
                results.append(ChunkResult(
                    content=chunk_text,
                    chunk_index=idx,
                    chunk_type="html",
                    metadata={"title": title, "paragraph_index": para_idx},
                ))
                idx += 1

    return results
